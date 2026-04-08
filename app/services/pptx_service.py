from collections.abc import Sequence
from io import BytesIO

from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.models.render_spec import ResolvedBlock, ResolvedSlide
from app.models.schema import Presentation
from app.services.layout_resolver import LayoutResolver
from app.services.slide_format import (
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    STANDARD_WIDESCREEN_NAME,
)
from app.utils.slide_design import is_japanese_text


class PptxRenderError(Exception):
    pass


class PptxService:
    SLIDE_WIDTH = SLIDE_WIDTH
    SLIDE_HEIGHT = SLIDE_HEIGHT
    FONT_FAMILY = "Meiryo"
    TITLE_COLOR = RGBColor(15, 23, 42)
    BODY_COLOR = RGBColor(51, 65, 85)
    MUTED_COLOR = RGBColor(71, 85, 105)
    ACCENT_COLOR = RGBColor(14, 116, 144)
    ACCENT_SOFT = RGBColor(224, 242, 254)
    BORDER_COLOR = RGBColor(203, 213, 225)
    BACKGROUND_COLOR = RGBColor(242, 247, 252)
    PANEL_COLOR = RGBColor(255, 255, 255)
    PANEL_ALT_COLOR = RGBColor(248, 250, 252)
    SUMMARY_PANEL_COLOR = RGBColor(236, 246, 250)
    TITLE_BAND_COLOR = RGBColor(15, 118, 110)
    TITLE_BAND_SOFT = RGBColor(204, 251, 241)
    PRIMARY_BLUE = RGBColor(37, 99, 235)
    SOFT_BLUE = RGBColor(219, 234, 254)
    LIGHT_BLUE = RGBColor(239, 246, 255)

    def __init__(self) -> None:
        self.layout_resolver = LayoutResolver()

    def render_pptx(self, presentation: Presentation) -> bytes:
        resolved = self.layout_resolver.resolve_presentation(presentation)
        deck = PptxPresentation()
        deck.slide_width = self.SLIDE_WIDTH
        deck.slide_height = self.SLIDE_HEIGHT
        deck.core_properties.title = presentation.deck_title
        deck.core_properties.subject = f"Generated slide deck ({STANDARD_WIDESCREEN_NAME})"
        deck.core_properties.language = self._infer_document_language(presentation)

        for slide_index, slide in enumerate(resolved.slides, start=1):
            pptx_slide = deck.slides.add_slide(deck.slide_layouts[6])
            self._apply_background(pptx_slide)
            self._render_slide(pptx_slide, slide, slide_index=slide_index)

        output = BytesIO()
        deck.save(output)
        return output.getvalue()

    def _infer_document_language(self, presentation: Presentation) -> str:
        text = " ".join(
            [presentation.deck_title]
            + [slide.title for slide in presentation.slides]
            + [bullet for slide in presentation.slides for bullet in slide.bullets]
        )
        if is_japanese_text(text):
            return "ja-JP"
        return "en-US"

    def _is_japanese_blocks(self, blocks: Sequence[ResolvedBlock]) -> bool:
        text = " ".join(
            [block.heading or "" for block in blocks]
            + [item.text for block in blocks for item in block.items]
        )
        return is_japanese_text(text)

    def _render_slide(self, pptx_slide, slide: ResolvedSlide, *, slide_index: int) -> None:
        if slide.pattern == "title_hero":
            self._render_title_slide(pptx_slide, slide)
        elif slide.pattern == "agenda_steps":
            self._render_agenda_slide(pptx_slide, slide)
        elif slide.pattern == "split_columns":
            self._render_two_column_slide(pptx_slide, slide)
        elif slide.pattern == "action_summary":
            self._render_summary_slide(pptx_slide, slide)
        elif slide.pattern == "standard_list":
            self._render_linear_content_slide(pptx_slide, slide)
        else:
            self._render_standard_slide(pptx_slide, slide)

    def _apply_background(self, pptx_slide) -> None:
        fill = pptx_slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.BACKGROUND_COLOR

        canvas = pptx_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.45),
            Inches(0.42),
            Inches(12.18),
            Inches(6.38),
        )
        self._set_shape_fill(canvas, RGBColor(255, 255, 255))
        canvas.line.color.rgb = RGBColor(219, 234, 254)

    def _render_title_slide(self, pptx_slide, slide: ResolvedSlide) -> None:
        title_box = pptx_slide.shapes.add_textbox(
            Inches(1.15),
            Inches(2.15),
            Inches(11.0),
            Inches(1.15),
        )
        self._configure_text_frame(title_box.text_frame, vertical_anchor=MSO_ANCHOR.MIDDLE)
        self._add_paragraph(
            title_box.text_frame,
            slide.title,
            font_size=Pt(30),
            bold=True,
            color=self.PRIMARY_BLUE,
            alignment=PP_ALIGN.CENTER,
        )

        summary_panel = pptx_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(1.3),
            Inches(3.95),
            Inches(10.75),
            Inches(1.45),
        )
        self._set_shape_fill(summary_panel, RGBColor(248, 250, 252))
        summary_panel.line.color.rgb = RGBColor(219, 234, 254)

        subtitle_box = pptx_slide.shapes.add_textbox(
            Inches(1.68),
            Inches(4.31),
            Inches(9.98),
            Inches(0.98),
        )
        self._configure_text_frame(subtitle_box.text_frame, vertical_anchor=MSO_ANCHOR.MIDDLE)
        for block in slide.blocks[:1]:
            for item in block.items[:3]:
                self._add_paragraph(
                    subtitle_box.text_frame,
                    item.text,
                    font_size=Pt(17),
                    color=self.BODY_COLOR,
                    alignment=PP_ALIGN.CENTER,
                )

    def _render_agenda_slide(self, pptx_slide, slide: ResolvedSlide) -> None:
        self._add_plain_title(pptx_slide, slide.title)
        self._render_block_cards(pptx_slide, slide.blocks)

    def _render_standard_slide(self, pptx_slide, slide: ResolvedSlide) -> None:
        self._add_plain_title(pptx_slide, slide.title)
        self._render_block_cards(pptx_slide, slide.blocks)

    def _render_linear_content_slide(self, pptx_slide, slide: ResolvedSlide) -> None:
        self._add_plain_title(pptx_slide, slide.title)

        panel = pptx_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(1.04),
            Inches(2.24),
            Inches(10.34),
            Inches(3.42),
        )
        self._set_shape_fill(panel, self.PANEL_ALT_COLOR)
        panel.line.color.rgb = self.BORDER_COLOR

        body_box = pptx_slide.shapes.add_textbox(
            Inches(1.42),
            Inches(2.66),
            Inches(9.58),
            Inches(2.52),
        )
        self._configure_text_frame(body_box.text_frame)
        for block in slide.blocks:
            for item in block.items:
                self._add_paragraph(
                    body_box.text_frame,
                    item.text,
                    font_size=Pt(20),
                    color=self.BODY_COLOR,
                    bullet=True,
                )

    def _render_two_column_slide(self, pptx_slide, slide: ResolvedSlide) -> None:
        self._add_plain_title(pptx_slide, slide.title)

        left_panel = pptx_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(1.0),
            Inches(2.52),
            Inches(5.18),
            Inches(3.32),
        )
        right_panel = pptx_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(6.5),
            Inches(2.52),
            Inches(5.18),
            Inches(3.32),
        )
        for panel in (left_panel, right_panel):
            self._set_shape_fill(panel, self.PANEL_ALT_COLOR)
            panel.line.color.rgb = self.BORDER_COLOR

        left_heading = slide.blocks[0].heading or "左"
        right_heading = slide.blocks[1].heading or "右"
        self._add_panel_heading(pptx_slide, left_heading, left=Inches(1.28), top=Inches(2.78))
        self._add_panel_heading(pptx_slide, right_heading, left=Inches(6.78), top=Inches(2.78))

        left_bullets = [item.text for item in slide.blocks[0].items]
        right_bullets = [item.text for item in slide.blocks[1].items]
        self._render_panel_bullets(pptx_slide, left_bullets, left=Inches(1.28), top=Inches(3.18))
        self._render_panel_bullets(pptx_slide, right_bullets, left=Inches(6.78), top=Inches(3.18))

    def _render_summary_slide(self, pptx_slide, slide: ResolvedSlide) -> None:
        self._add_plain_title(pptx_slide, slide.title)
        positions = [(Inches(1.02), Inches(2.7)), (Inches(6.55), Inches(2.7))]
        for (left, top), block in zip(positions, slide.blocks, strict=False):
            card = pptx_slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                left,
                top,
                Inches(4.78),
                Inches(2.45),
            )
            self._set_shape_fill(card, self.SOFT_BLUE)
            card.line.color.rgb = self.PRIMARY_BLUE

            header_box = pptx_slide.shapes.add_textbox(
                left + Inches(0.25),
                top + Inches(0.22),
                Inches(2.6),
                Inches(0.35),
            )
            self._configure_text_frame(header_box.text_frame, vertical_anchor=MSO_ANCHOR.MIDDLE)
            self._add_paragraph(
                header_box.text_frame,
                block.heading or ("アクション" if self._is_japanese_blocks(slide.blocks) else "Action"),
                font_size=Pt(13),
                bold=True,
                color=self.PRIMARY_BLUE,
            )

            body_box = pptx_slide.shapes.add_textbox(
                left + Inches(0.25),
                top + Inches(0.72),
                Inches(4.15),
                Inches(1.35),
            )
            self._configure_text_frame(body_box.text_frame)
            self._add_paragraph(
                body_box.text_frame,
                self._block_body(block),
                font_size=Pt(19),
                color=self.BODY_COLOR,
            )

    def _add_number_title(
        self,
        pptx_slide,
        title: str,
        slide_index: int | None,
        *,
        role_label: str | None = None,
    ) -> None:
        if slide_index is not None:
            number_text = f"{slide_index:02d}"
        else:
            number_text = ""
        if number_text:
            badge = pptx_slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(1.2),
                Inches(0.9),
                Inches(0.86),
                Inches(0.58),
            )
            self._set_shape_fill(badge, self.PRIMARY_BLUE)
            badge.line.fill.background()
            badge_text = badge.text_frame
            self._configure_text_frame(badge_text, vertical_anchor=MSO_ANCHOR.MIDDLE)
            self._add_paragraph(
                badge_text,
                number_text,
                font_size=Pt(16),
                bold=True,
                color=RGBColor(255, 255, 255),
                alignment=PP_ALIGN.CENTER,
            )
            title_left = Inches(2.16)
            title_width = Inches(9.5)
        else:
            title_left = Inches(1.18)
            title_width = Inches(10.5)

        title_box = pptx_slide.shapes.add_textbox(
            title_left,
            Inches(0.85),
            title_width,
            Inches(0.7),
        )
        self._configure_text_frame(title_box.text_frame, vertical_anchor=MSO_ANCHOR.MIDDLE)
        self._add_paragraph(
            title_box.text_frame,
            title,
            font_size=Pt(25),
            bold=True,
            color=self.PRIMARY_BLUE,
        )
        if role_label:
            self._add_chip(
                pptx_slide,
                role_label,
                left=title_left,
                top=Inches(1.64),
                width=Inches(1.85),
                fill_color=self.ACCENT_SOFT,
                line_color=self.ACCENT_COLOR,
                text_color=self.ACCENT_COLOR,
            )

    def _add_plain_title(self, pptx_slide, title: str) -> None:
        title_box = pptx_slide.shapes.add_textbox(
            Inches(1.18),
            Inches(0.85),
            Inches(10.5),
            Inches(0.7),
        )
        self._configure_text_frame(title_box.text_frame, vertical_anchor=MSO_ANCHOR.MIDDLE)
        self._add_paragraph(
            title_box.text_frame,
            title,
            font_size=Pt(25),
            bold=True,
            color=self.PRIMARY_BLUE,
        )

    def _render_block_cards(self, pptx_slide, blocks: tuple[ResolvedBlock, ...]) -> None:
        positions = self._content_card_positions(max(len(blocks), 1))
        for block, (left, top, width, height) in zip(blocks, positions, strict=False):
            self._add_content_card(
                pptx_slide,
                block.heading or ("要点" if self._is_japanese_blocks(blocks) else "Point"),
                self._block_body(block),
                left=left,
                top=top,
                width=width,
                height=height,
            )

    def _add_content_card(self, pptx_slide, label: str, body: str, *, left, top, width, height):
        card = pptx_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            top,
            width,
            height,
        )
        self._set_shape_fill(card, self.SOFT_BLUE)
        card.line.color.rgb = RGBColor(191, 219, 254)

        label_box = pptx_slide.shapes.add_textbox(
            left + Inches(0.22),
            top + Inches(0.18),
            width - Inches(0.44),
            Inches(0.3),
        )
        self._configure_text_frame(label_box.text_frame, vertical_anchor=MSO_ANCHOR.MIDDLE)
        self._add_paragraph(
            label_box.text_frame,
            label,
            font_size=Pt(13),
            bold=True,
            color=self.PRIMARY_BLUE,
        )

        body_box = pptx_slide.shapes.add_textbox(
            left + Inches(0.22),
            top + Inches(0.78),
            width - Inches(0.44),
            height - Inches(1.0),
        )
        self._configure_text_frame(body_box.text_frame)
        for line in body.split("<br>"):
            self._add_paragraph(
                body_box.text_frame,
                line,
                font_size=Pt(20),
                bold=False,
                color=self.BODY_COLOR,
            )
        return card

    def _content_card_positions(self, bullet_count: int) -> list[tuple[object, object, object, object]]:
        if bullet_count <= 1:
            return [(Inches(1.08), Inches(2.28), Inches(10.2), Inches(2.88))]
        if bullet_count == 2:
            return [
                (Inches(1.05), Inches(2.28), Inches(4.82), Inches(2.92)),
                (Inches(6.44), Inches(2.28), Inches(4.82), Inches(2.92)),
            ]
        if bullet_count == 3:
            return [
                (Inches(0.92), Inches(2.32), Inches(3.44), Inches(3.0)),
                (Inches(4.94), Inches(2.32), Inches(3.44), Inches(3.0)),
                (Inches(8.96), Inches(2.32), Inches(3.44), Inches(3.0)),
            ]
        return [
            (Inches(1.05), Inches(2.12), Inches(4.82), Inches(1.56)),
            (Inches(6.44), Inches(2.12), Inches(4.82), Inches(1.56)),
            (Inches(1.05), Inches(4.05), Inches(4.82), Inches(1.56)),
            (Inches(6.44), Inches(4.05), Inches(4.82), Inches(1.56)),
        ]

    def _add_arrow_connector(self, pptx_slide, *, left, top, width) -> None:
        arrow = pptx_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            left,
            top,
            width,
            Inches(0.55),
        )
        self._set_shape_fill(arrow, RGBColor(96, 165, 250))
        arrow.line.fill.background()

    def _render_panel_bullets(self, pptx_slide, bullets: Sequence[str], *, left, top) -> None:
        body_box = pptx_slide.shapes.add_textbox(left, top, Inches(4.58), Inches(2.25))
        self._configure_text_frame(body_box.text_frame)
        if not bullets:
            self._add_paragraph(
                body_box.text_frame,
                "補足事項なし" if is_japanese_text(" ".join(bullets)) else "No additional notes",
                font_size=Pt(16),
                color=self.MUTED_COLOR,
            )
            return
        for bullet in bullets:
            self._add_paragraph(
                body_box.text_frame,
                bullet,
                font_size=Pt(16),
                color=self.BODY_COLOR,
                bullet=True,
            )

    def _add_panel_heading(self, pptx_slide, text: str, *, left, top) -> None:
        heading_box = pptx_slide.shapes.add_textbox(left, top, Inches(1.5), Inches(0.3))
        self._configure_text_frame(heading_box.text_frame, vertical_anchor=MSO_ANCHOR.MIDDLE)
        self._add_paragraph(
            heading_box.text_frame,
            text,
            font_size=Pt(12),
            bold=True,
            color=self.ACCENT_COLOR,
        )

    def _add_chip(
        self,
        pptx_slide,
        text: str,
        *,
        left,
        top,
        width,
        fill_color: RGBColor | None = None,
        line_color: RGBColor | None = None,
        text_color: RGBColor | None = None,
    ) -> None:
        chip = pptx_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            top,
            width,
            Inches(0.4),
        )
        self._set_shape_fill(chip, fill_color or self.ACCENT_SOFT)
        chip.line.color.rgb = line_color or self.ACCENT_COLOR
        chip_text = chip.text_frame
        self._configure_text_frame(chip_text, vertical_anchor=MSO_ANCHOR.MIDDLE)
        self._add_paragraph(
            chip_text,
            text,
            font_size=Pt(10),
            bold=True,
            color=text_color or self.ACCENT_COLOR,
            alignment=PP_ALIGN.CENTER,
        )

    def _block_body(self, block: ResolvedBlock) -> str:
        return "<br>".join(item.text for item in block.items if item.text.strip())

    def _configure_text_frame(self, text_frame, *, vertical_anchor=MSO_ANCHOR.TOP) -> None:
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.vertical_anchor = vertical_anchor

    def _add_paragraph(
        self,
        text_frame,
        text: str,
        *,
        font_size,
        color: RGBColor,
        bold: bool = False,
        bullet: bool = False,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
    ) -> None:
        paragraph = text_frame.paragraphs[0] if not text_frame.text else text_frame.add_paragraph()
        paragraph.text = text
        paragraph.alignment = alignment
        paragraph.level = 0
        paragraph.font.size = font_size
        paragraph.font.bold = bold
        paragraph.font.name = self.FONT_FAMILY
        paragraph.font.color.rgb = color
        paragraph.bullet = bullet
        paragraph.line_spacing = 1.25
        if bullet:
            paragraph.space_after = Pt(7)
        else:
            paragraph.space_after = Pt(4)

    def _set_shape_fill(self, shape, color: RGBColor) -> None:
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = color
