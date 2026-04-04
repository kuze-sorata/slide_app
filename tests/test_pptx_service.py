from io import BytesIO

from pptx import Presentation as PptxPresentation

from app.models.schema import Presentation
from app.services.pptx_service import PptxService


def sample_presentation() -> Presentation:
    return Presentation.model_validate(
        {
            "deck_title": "営業進捗の共有",
            "slides": [
                {
                    "id": "slide-1",
                    "type": "title",
                    "title": "営業進捗の共有",
                    "bullets": ["部長向け共有", "現状と課題を提示"],
                    "layout": "layout1",
                },
                {
                    "id": "slide-2",
                    "type": "content",
                    "title": "重点案件の進捗",
                    "bullets": ["A案件: 契約間近", "B案件: 交渉停滞中"],
                    "layout": "layout2",
                },
                {
                    "id": "slide-3",
                    "type": "summary",
                    "title": "今後のアクション",
                    "bullets": ["B案件の戦略見直し", "週次レビューの強化"],
                    "layout": "layout4",
                },
            ],
        }
    )


def test_render_pptx_outputs_valid_file() -> None:
    service = PptxService()

    pptx_bytes = service.render_pptx(sample_presentation())

    assert pptx_bytes.startswith(b"PK")
    deck = PptxPresentation(BytesIO(pptx_bytes))
    assert len(deck.slides) == 3
    assert deck.core_properties.title == "営業進捗の共有"

    slide_texts = [
        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        for slide in deck.slides
    ]

    assert "営業進捗の共有" in slide_texts[0]
    assert "社内向けスライド草案" not in slide_texts[0]
    assert "部長向け共有" in slide_texts[0]
    assert "資料の主題" not in slide_texts[0]
    assert "社内向け草案" not in slide_texts[0]
    assert "重点案件の進捗" in slide_texts[1]
    assert "A案件: 契約間近" in slide_texts[1]
    assert "今後のアクション" in slide_texts[2]
    assert "週次レビューの強化" in slide_texts[2]
    assert "アクション 1" in slide_texts[2] or "アクション 2" in slide_texts[2]


def test_render_pptx_layout3_splits_bullets_across_columns() -> None:
    service = PptxService()
    presentation = Presentation.model_validate(
        {
            "deck_title": "案件整理",
            "slides": [
                {
                    "id": "slide-1",
                    "type": "title",
                    "title": "案件整理",
                    "bullets": ["概要共有"],
                    "layout": "layout1",
                },
                {
                    "id": "slide-2",
                    "type": "content",
                    "title": "対応方針",
                    "bullets": ["優先順位を見直す", "担当を再配置する", "レビュー頻度を上げる"],
                    "layout": "layout3",
                },
                {
                    "id": "slide-3",
                    "type": "summary",
                    "title": "進め方",
                    "bullets": ["当日中に整理する"],
                    "layout": "layout4",
                },
            ],
        }
    )

    deck = PptxPresentation(BytesIO(service.render_pptx(presentation)))
    content_slide_text = "\n".join(
        shape.text for shape in deck.slides[1].shapes if hasattr(shape, "text")
    )

    assert "対応方針" in content_slide_text
    assert "打ち手 1" in content_slide_text
    assert "打ち手 2" in content_slide_text
    assert "優先順位を見直す" in content_slide_text
    assert "担当を再配置する" in content_slide_text
    assert "レビュー頻度を上げる" in content_slide_text


def test_render_pptx_layout2_keeps_three_parallel_points_separate() -> None:
    service = PptxService()
    presentation = Presentation.model_validate(
        {
            "deck_title": "サービス整理",
            "slides": [
                {
                    "id": "slide-1",
                    "type": "title",
                    "title": "サービス整理",
                    "bullets": ["営業部向け共有"],
                    "layout": "layout1",
                },
                {
                    "id": "slide-2",
                    "type": "content",
                    "title": "主要サービスを整理する",
                    "bullets": ["既存顧客向け支援", "導入初期の伴走支援", "運用改善の定着支援"],
                    "layout": "layout2",
                },
                {
                    "id": "slide-3",
                    "type": "summary",
                    "title": "訴求順を決める",
                    "bullets": ["3サービスの見せ方を統一する"],
                    "layout": "layout4",
                },
            ],
        }
    )

    deck = PptxPresentation(BytesIO(service.render_pptx(presentation)))
    content_slide_text = "\n".join(
        shape.text for shape in deck.slides[1].shapes if hasattr(shape, "text")
    )

    assert "要点 1" in content_slide_text
    assert "要点 2" in content_slide_text
    assert "要点 3" in content_slide_text
    assert "既存顧客向け支援" in content_slide_text
    assert "導入初期の伴走支援" in content_slide_text
    assert "運用改善の定着支援" in content_slide_text


def test_render_pptx_layout2_keeps_non_parallel_points_as_vertical_bullets() -> None:
    service = PptxService()
    presentation = Presentation.model_validate(
        {
            "deck_title": "進捗共有",
            "slides": [
                {
                    "id": "slide-1",
                    "type": "title",
                    "title": "進捗共有",
                    "bullets": ["営業部向け共有"],
                    "layout": "layout1",
                },
                {
                    "id": "slide-2",
                    "type": "content",
                    "title": "全体進捗は計画を下回る",
                    "bullets": ["主要案件の受注時期が遅延", "新規リード獲得が目標未達"],
                    "layout": "layout2",
                },
                {
                    "id": "slide-3",
                    "type": "summary",
                    "title": "対応方針",
                    "bullets": ["優先案件を見直す"],
                    "layout": "layout4",
                },
            ],
        }
    )

    deck = PptxPresentation(BytesIO(service.render_pptx(presentation)))
    content_slide_text = "\n".join(
        shape.text for shape in deck.slides[1].shapes if hasattr(shape, "text")
    )

    assert "全体進捗は計画を下回る" in content_slide_text
    assert "主要案件の受注時期が遅延" in content_slide_text
    assert "新規リード獲得が目標未達" in content_slide_text
    assert "要点 1" not in content_slide_text
